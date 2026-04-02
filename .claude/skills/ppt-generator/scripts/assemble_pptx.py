#!/usr/bin/env python3
"""
PPTX 组装脚本
将生成的幻灯片图片组装成可编辑的 PowerPoint 文件
"""
import argparse
import yaml
from pathlib import Path
from typing import Optional, List, Dict

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️ python-pptx 未安装，请运行: pip install python-pptx")


def load_config(ppt_dir: Path) -> Optional[Dict]:
    """加载 PPT 配置"""
    config_file = ppt_dir / 'ppt_content.yaml'
    if config_file.exists():
        with open(config_file, 'r', encoding='utf-8') as f:
            return yaml.safe_load(f)
    return None


def get_slide_images(ppt_dir: Path) -> List[Path]:
    """获取所有幻灯片图片，按页码排序"""
    images = list(ppt_dir.glob('slide_*.png'))
    images.sort(key=lambda x: int(x.stem.split('_')[1]))
    return images


def assemble_pptx(
    input_dir: str,
    output_path: str,
    add_text: bool = True
) -> str:
    """
    组装 PPTX 文件
    
    Args:
        input_dir: 包含 slide_XX.png 的目录
        output_path: 输出 PPTX 路径
        add_text: 是否添加文字层（基于 ppt_content.yaml）
    
    Returns:
        输出文件路径
    """
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx 未安装")
    
    ppt_dir = Path(input_dir)
    output_file = Path(output_path)
    
    # 获取幻灯片图片
    slide_images = get_slide_images(ppt_dir)
    if not slide_images:
        raise FileNotFoundError(f"未找到幻灯片图片: {ppt_dir}")
    
    print(f"📦 组装 PPTX: {len(slide_images)} 张幻灯片")
    
    # 加载配置
    config = load_config(ppt_dir)
    slides_config = config.get('slides', []) if config else []
    
    # 创建 PPT
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    
    # 使用空白布局
    blank_layout = prs.slide_layouts[6]  # 空白布局
    
    for idx, img_path in enumerate(slide_images):
        print(f"   添加幻灯片 #{idx + 1}: {img_path.name}")
        
        # 添加幻灯片
        slide = prs.slides.add_slide(blank_layout)
        
        # 添加背景图片（全屏）
        slide.shapes.add_picture(
            str(img_path),
            Inches(0), Inches(0),
            width=prs.slide_width,
            height=prs.slide_height
        )
        
        # 添加文字层（如果有配置）
        if add_text and idx < len(slides_config):
            slide_info = slides_config[idx]
            add_text_to_slide(slide, slide_info, prs)
    
    # 保存
    output_file.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_file))
    
    print(f"✅ PPTX 已保存: {output_file}")
    return str(output_file)


def add_text_to_slide(slide, slide_info: Dict, prs: Presentation):
    """
    向幻灯片添加文字层
    
    基于布局类型添加相应的文字框
    """
    layout = slide_info.get('layout', 'content_left')
    content = slide_info.get('content', '')
    
    if not content:
        return
    
    # 根据布局类型确定文字位置
    if layout == 'cover':
        # 封面：居中标题
        add_title_textbox(
            slide,
            text=content,
            left=Inches(1.5),
            top=Inches(3),
            width=Inches(10.333),
            height=Inches(1.5),
            font_size=Pt(54),
            bold=True,
            align=PP_ALIGN.CENTER
        )
    
    elif layout == 'section_divider':
        # 章节分隔：居中
        add_title_textbox(
            slide,
            text=content,
            left=Inches(2),
            top=Inches(3.2),
            width=Inches(9.333),
            height=Inches(1.2),
            font_size=Pt(44),
            bold=True,
            align=PP_ALIGN.CENTER
        )
    
    elif layout == 'content_left':
        # 左对齐内容
        lines = content.strip().split('\n')
        if lines:
            # 标题
            add_title_textbox(
                slide,
                text=lines[0].lstrip('•·- '),
                left=Inches(1.2),
                top=Inches(1),
                width=Inches(11),
                height=Inches(0.8),
                font_size=Pt(36),
                bold=True,
                align=PP_ALIGN.LEFT
            )
            # 内容
            if len(lines) > 1:
                body_text = '\n'.join(lines[1:])
                add_body_textbox(
                    slide,
                    text=body_text,
                    left=Inches(1.2),
                    top=Inches(2),
                    width=Inches(11),
                    height=Inches(5),
                    font_size=Pt(18),
                    align=PP_ALIGN.LEFT
                )
    
    elif layout == 'content_center':
        # 居中强调
        add_title_textbox(
            slide,
            text=content,
            left=Inches(2.5),
            top=Inches(2.5),
            width=Inches(8.333),
            height=Inches(3),
            font_size=Pt(28),
            bold=False,
            align=PP_ALIGN.CENTER
        )
    
    elif layout == 'content_two_column':
        # 双栏：简化处理，放在左侧
        add_body_textbox(
            slide,
            text=content,
            left=Inches(1.2),
            top=Inches(2),
            width=Inches(11),
            height=Inches(5),
            font_size=Pt(18),
            align=PP_ALIGN.LEFT
        )


def add_title_textbox(slide, text: str, left, top, width, height, 
                      font_size, bold: bool, align):
    """添加标题文字框"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = RgbColor(0x1D, 0x1D, 0x1F)  # Apple 黑色
    p.alignment = align


def add_body_textbox(slide, text: str, left, top, width, height,
                     font_size, align):
    """添加正文文字框"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    
    # 按行添加
    lines = text.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = line
        p.font.size = font_size
        p.font.color.rgb = RgbColor(0x1D, 0x1D, 0x1F)
        p.alignment = align
        p.space_after = Pt(8)


def main():
    parser = argparse.ArgumentParser(description='将幻灯片图片组装成 PPTX 文件')
    parser.add_argument('--input_dir', '-i', required=True, 
                        help='包含 slide_XX.png 的目录')
    parser.add_argument('--output', '-o', required=True,
                        help='输出 PPTX 文件路径')
    parser.add_argument('--no-text', action='store_true',
                        help='不添加文字层（仅图片）')
    
    args = parser.parse_args()
    
    try:
        output_path = assemble_pptx(
            input_dir=args.input_dir,
            output_path=args.output,
            add_text=not args.no_text
        )
        print(f"\n🎉 完成！文件: {output_path}")
    except Exception as e:
        print(f"\n❌ 错误: {e}")
        raise


if __name__ == '__main__':
    main()
